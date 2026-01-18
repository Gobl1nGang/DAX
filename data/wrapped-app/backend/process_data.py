import json
import glob
import os
import re
from collections import defaultdict, Counter
from datetime import datetime, timezone, timedelta


def fix_text(text):
    if not isinstance(text, str) or not text:
        return text
    try:
        # Facebook JSON often encodes UTF-8 bytes as Latin-1 characters.
        return text.encode('latin1').decode('utf-8')
    except (UnicodeEncodeError, UnicodeDecodeError):
        try:
            return text.encode('cp1252').decode('utf-8')
        except:
            return text

def is_url(text):
    if not isinstance(text, str):
        return False
    text_lower = text.lower()
    return text_lower.startswith('http') or text_lower.startswith('www') or '://' in text_lower

def process_uploaded_data(files_data):
    # Data structures for stats
    user_word_counts = defaultdict(Counter)
    total_word_counts = Counter()
    messages_with_likes = []
    quiz_pool = []
    user_image_counts = Counter()

    # Special Rankings Data
    all_messages = []  # Flattened list of all messages for time/reply analysis
    user_likes_received = Counter()
    user_messages_sent = Counter()
    user_reels_sent = Counter()
    user_longest_words = defaultdict(lambda: {"word": "", "length": 0})
    user_likes_given = Counter() # Most Messages Liked
    likes_matrix = defaultdict(lambda: defaultdict(int)) # actor -> sender likes
    user_total_chars = Counter()
    month_counts = Counter()
    
    # Combined profanity pattern for faster matching
    _PROFANITY_RE = re.compile(r'\b([fF]+[uU]+[cC]+[kK]+|[sS]+[hH]+[iI]+[tT]+|[bB]+[iI]+[tT]+[cC]+[hH]+|[dD]+[aA]+[mM]+[nN]+|[aA]+[sS]+[sS]+|[hH]+[eE]+[lL]+[lL]+)[a-zA-Z]*\b')
    user_profanity_counts = Counter()

    STOPWORDS = {
        'i', 'me', 'my', 'myself', 'we', 'our', 'ours', 'ourselves', 'you', "you're", "you've", "you'll", "you'd",
        'your', 'yours', 'yourself', 'yourselves', 'he', 'him', 'his', 'himself', 'she', "she's", 'her', 'hers',
        'herself', 'it', "it's", 'its', 'itself', 'they', 'them', 'their', 'theirs', 'themselves', 'what', 'which',
        'who', 'whom', 'this', 'that', "that'll", 'these', 'those', 'am', 'is', 'are', 'was', 'were', 'be', 'been',
        'being', 'have', 'has', 'had', 'having', 'do', 'does', 'did', 'doing', 'a', 'an', 'the', 'and', 'but', 'if',
        'or', 'because', 'as', 'until', 'while', 'of', 'at', 'by', 'for', 'with', 'about', 'against', 'between',
        'into', 'through', 'during', 'before', 'after', 'above', 'below', 'to', 'from', 'up', 'down', 'in', 'out',
        'on', 'off', 'over', 'under', 'again', 'further', 'then', 'once', 'here', 'there', 'when', 'where', 'why',
        'how', 'all', 'any', 'both', 'each', 'few', 'more', 'most', 'other', 'some', 'such', 'no', 'nor', 'not',
        'only', 'own', 'same', 'so', 'than', 'too', 'very', 's', 't', 'can', 'will', 'just', 'don', "don't", 'should',
        "should've", 'now', 'd', 'll', 'm', 'o', 're', 've', 'y', 'ain', 'aren', "aren't", 'couldn', "couldn't",
        'didn', "didn't", 'doesn', "doesn't", 'hadn', "hadn't", 'hasn', "hasn't", 'haven', "haven't", 'isn', "isn't",
        'ma', 'mightn', "mightn't", 'mustn', "mustn't", 'needn', "needn't", 'shan', "shan't", 'shouldn', "shouldn't",
        'wasn', "wasn't", 'weren', "weren't", 'won', "won't", 'wouldn', "wouldn't", 'like', 'get', 'got', 'know',
        'think', 'going', 'really', 'yeah', 'lol', 'lmao', 'good', 'well', 'much', 'see', 'want', 'one', 'even',
        'message', 'liked', 'im', 'guys', 'sent', 'attachment', 'dont', 'ur', 'thats',
        'yeah', 'yes', 'no', 'ok', 'okay', 'like', 'get', 'got', 'know', 'think', 'going',
        'really', 'good', 'well', 'much', 'see', 'want', 'one', 'even', 'right', 'mean',
        'thing', 'things', 'still', 'actually', 'probably', 'maybe', 'sure', 'back',
        'come', 'take', 'look', 'want', 'gonna', 'wanna', 'gotta', 'lol', 'lmao', 'lmfao',
        'haha', 'hahaha', 'ah', 'oh', 'hey', 'hi', 'hello'
    }

    # Small mapping to expand common contractions so we don't end up with tokens like "thats"
    CONTRACTIONS_MAP = {
        "ain't": "is not",
        "aren't": "are not",
        "can't": "cannot",
        "can't've": "cannot have",
        "could've": "could have",
        "couldn't": "could not",
        "couldn't've": "could not have",
        "didn't": "did not",
        "doesn't": "does not",
        "don't": "do not",
        "hadn't": "had not",
        "hadn't've": "had not have",
        "hasn't": "has not",
        "haven't": "have not",
        "he'd": "he would",
        "he'd've": "he would have",
        "he'll": "he will",
        "he's": "he is",
        "i'd": "i would",
        "i'd've": "i would have",
        "i'll": "i will",
        "i'm": "i am",
        "i've": "i have",
        "isn't": "is not",
        "it's": "it is",
        "let's": "let us",
        "mightn't": "might not",
        "mustn't": "must not",
        "shan't": "shall not",
        "she's": "she is",
        "should've": "should have",
        "shouldn't": "should not",
        "that's": "that is",
        "there's": "there is",
        "they're": "they are",
        "they've": "they have",
        "we're": "we are",
        "we've": "we have",
        "weren't": "were not",
        "what's": "what is",
        "won't": "will not",
        "wouldn't": "would not",
        "you'd": "you would",
        "you'll": "you will",
        "you're": "you are",
        # Add defensive lowercase variants without apostrophes commonly produced by earlier cleaning
        'thats': 'that is',
        'dont': 'do not',
        'im': 'i am',
        'ive': 'i have',
        'cant': 'cannot',
        'its': 'it is'
    }

    # Pre-compile contractions regex for speed
    _CONTRACTIONS_RE = re.compile(r'\b(' + '|'.join(re.escape(key) for key in CONTRACTIONS_MAP.keys()) + r')\b', re.IGNORECASE)

    def expand_contractions(text):
        if not isinstance(text, str) or not text:
            return text
        # Normalize simple unicode apostrophes
        t = text.replace('\u2019', "'").replace('\u2018', "'")
        return _CONTRACTIONS_RE.sub(lambda m: CONTRACTIONS_MAP.get(m.group(0).lower(), m.group(0)), t.lower())

    # Fast word filter heuristic
    def _fast_filter_words(words):
        allowed = []
        for w in words:
            wl = w.lower()
            if len(wl) <= 2: continue
            if wl in STOPWORDS: continue
            # Simple heuristic: ignore common pronouns/determiners that might have slipped through
            if wl in {'that', 'this', 'these', 'those', 'they', 'them', 'their', 'there'}: continue
            if re.match(r'^[a-zA-Z]+$', wl):
                allowed.append(wl)
        return allowed

    collected_media_basenames = set()

    for file_info in files_data:
        try:
            data = json.loads(file_info['content'])
            
            if 'messages' not in data:
                continue

            for message in data['messages']:
                if 'sender_name' not in message:
                    continue
                    
                sender = fix_text(message['sender_name'])
                
                # Exclude Meta AI
                if sender == "Meta AI":
                    continue
                    
                user_messages_sent[sender] += 1

                # Store for chronological processing
                if 'timestamp_ms' in message:
                    # Fix sender name in stored message too for consistency later
                    message['sender_name'] = sender
                    all_messages.append(message)

                # 1. Word Counts & Longest Word & Profanity
                if 'content' in message:
                    content = fix_text(message['content'])
                    # Expand contractions (so "that's" -> "that is" and gets removed by stopwords)
                    content = expand_contractions(content)
                    user_total_chars[sender] += len(content)

                    # Split into words first
                    raw_words = content.split()

                    # Filter out URLs
                    non_url_words = [w for w in raw_words if not is_url(w)]

                    # Process for Word Counts
                    processed_words = []
                    for w in non_url_words:
                        # Check for profanity in a single regex pass
                        if _PROFANITY_RE.match(w):
                            user_profanity_counts[sender] += 1

                        # Basic tokenization: lowercase and remove non-alphanumeric characters
                        cleaned = re.sub(r'[^\w\s]', '', w.lower())
                        # Also strip surrounding punctuation
                        cleaned = cleaned.strip()
                        if cleaned and cleaned not in STOPWORDS and len(cleaned) > 1:
                            processed_words.append(cleaned)

                    # Filter words using fast heuristic
                    pos_words = _fast_filter_words(processed_words)
                    user_word_counts[sender].update(pos_words)
                    total_word_counts.update(pos_words)

                    # Track longest word (Exclude links and non-alpha)
                    for word in non_url_words:
                        # Clean word for length check (remove punctuation from ends)
                        clean_word = re.sub(r'^[^\w]+|[^\w]+$', '', word)

                        # Strict check: Must be purely English letters to avoid emojis/symbols/mojibake
                        if not re.match(r'^[a-zA-Z]+$', clean_word):
                            continue

                        if len(clean_word) > user_longest_words[sender]["length"]:
                            user_longest_words[sender] = {"word": clean_word, "length": len(clean_word)}


                    # Check for Reels
                    if "instagram.com/reel" in content:
                        user_reels_sent[sender] += 1
                    elif 'share' in message and isinstance(message.get('share'), dict) and 'link' in message['share'] and "instagram.com/reel" in message['share']['link']:
                        user_reels_sent[sender] += 1


                # 3. Most Liked Messages & Aura (Likes Received) & Likes Given & Matrix
                if 'reactions' in message:
                    like_count = len(message['reactions'])
                    if like_count > 0:
                        user_likes_received[sender] += like_count
                        
                        # Track who GAVE the likes
                        for reaction in message['reactions']:
                            if 'actor' in reaction:
                                actor = fix_text(reaction['actor'])
                                if actor == "Meta AI":
                                    continue
                                user_likes_given[actor] += 1
                                likes_matrix[actor][sender] += 1
                        
                        # Fix content encoding and filter out media-only messages
                        raw_content = message.get('content', '')
                        if not raw_content or raw_content == '[Media/No Content]':
                            continue
                            
                        fixed_content = fix_text(raw_content)
                        
                        msg_info = {
                            "content": fixed_content,
                            "sender": sender,
                            "likes": like_count,
                            "timestamp": message.get('timestamp_ms', 0)
                        }
                        
                        # 6. Quiz Pool (Sentences)
                        if 30 < len(fixed_content) < 150 and not is_url(fixed_content):
                            quiz_pool.append({"content": fixed_content, "sender": sender})

                        messages_with_likes.append(msg_info)

                # 6. Image/Media Stats (robust detection for screenshots and attachments)
                media_count = 0
                media_uris = []

                # Explicit photo lists
                photos = message.get('photos')
                if isinstance(photos, list):
                    media_count += len(photos)
                    for p in photos:
                        if isinstance(p, dict):
                            uri = p.get('uri') or p.get('uri_image') or p.get('uri_raw') or p.get('url')
                            if isinstance(uri, str) and uri:
                                media_uris.append(uri)
                elif photos:
                    media_count += 1

                # Videos
                videos = message.get('videos')
                if isinstance(videos, list):
                    media_count += len(videos)
                    for v in videos:
                        if isinstance(v, dict):
                            uri = v.get('uri') or v.get('url')
                            if isinstance(uri, str) and uri:
                                media_uris.append(uri)
                elif videos:
                    media_count += 1

                # GIFs
                gifs = message.get('gifs')
                if isinstance(gifs, list):
                    media_count += len(gifs)
                    for g in gifs:
                        if isinstance(g, dict):
                            uri = g.get('uri') or g.get('url')
                            if isinstance(uri, str) and uri:
                                media_uris.append(uri)
                elif gifs:
                    media_count += 1

                # Sticker
                if 'sticker' in message:
                    media_count += 1
                    st = message.get('sticker')
                    if isinstance(st, dict):
                        uri = st.get('uri') or st.get('url')
                        if isinstance(uri, str) and uri:
                            media_uris.append(uri)

                # Attachments: often where screenshots live
                for att in message.get('attachments', []) or []:
                    if not isinstance(att, dict):
                        continue
                    atype = (att.get('type') or '').lower()

                    # Common attachment types that represent media
                    if atype in ('photo', 'image', 'video', 'gif', 'sticker', 'file'):
                        media_count += 1
                        media = att.get('media') or att.get('image') or att.get('photo') or {}
                        if isinstance(media, dict):
                            uri = media.get('uri') or media.get('uri_image') or media.get('uri_raw') or media.get('url') or media.get('src')
                            if isinstance(uri, str) and uri:
                                media_uris.append(uri)
                        continue

                    # Fallback: check nested media/image dicts for file URIs
                    media = att.get('media') or att.get('image') or att.get('photo') or {}
                    if isinstance(media, dict):
                        uri = media.get('uri') or media.get('uri_image') or media.get('url') or media.get('src')
                        if isinstance(uri, str) and re.search(r'\.(png|jpe?g|gif|bmp|webp)$', uri, flags=re.I):
                            media_count += 1
                            media_uris.append(uri)

                # Files array (some exports put attachments here)
                for f in message.get('files', []) or []:
                    if isinstance(f, dict):
                        uri = f.get('uri') or f.get('uri_raw') or f.get('name') or f.get('file')
                        if isinstance(uri, str) and re.search(r'\.(png|jpe?g|gif|bmp|webp)$', uri, flags=re.I):
                            media_count += 1
                            media_uris.append(uri)

                # Share links that might point to images
                share = message.get('share') or {}
                if isinstance(share, dict):
                    link = share.get('link') or share.get('href')
                    if isinstance(link, str) and re.search(r'\.(png|jpe?g|gif|bmp|webp)$', link, flags=re.I):
                        media_count += 1
                        media_uris.append(link)

                # Heuristic: filenames embedded in content (IMG_, Screenshot, date patterns)
                if media_count == 0:
                    content_for_media = message.get('content', '')
                    if isinstance(content_for_media, str) and re.search(r'(IMG[_-]\d+|Screenshot_?\d+|\d{8}_\d{6}|\.(png|jpe?g))', content_for_media, flags=re.I):
                        media_count += 1

                if media_count:
                    user_image_counts[sender] += media_count
                    # Expose detected URIs to the message for downstream zipping/collecting logic
                    if media_uris:
                        # Normalize URIs and collect basenames so the zipper can find files in the export
                        normalized = []
                        basenames = []
                        for uri in media_uris:
                            if not isinstance(uri, str):
                                continue
                            u = uri.strip()
                            # Keep data-URI or absolute URLs as-is
                            if u.startswith('data:') or re.search(r'https?://', u):
                                if u not in normalized:
                                    normalized.append(u)
                                continue

                            # Strip common leading path fragments
                            u_clean = re.sub(r'^(?:\./|\.\./|/)+', '', u)
                            # Some exports prefix with 'photos/' or 'images/' - keep full path but also capture basename
                            base = os.path.basename(u_clean)
                            if base:
                                if base not in basenames:
                                    basenames.append(base)
                                if u_clean not in normalized:
                                    normalized.append(u_clean)
                            else:
                                if u_clean not in normalized:
                                    normalized.append(u_clean)
                        
                        message['media_uris'] = normalized
                        message['media_basenames'] = basenames
                        for b in basenames:
                            collected_media_basenames.add(b)

        except Exception as e:
            print(f"Error processing file: {e}")

    # --- Compile Stats ---

    # Sort all messages by timestamp for time-based analysis
    all_messages.sort(key=lambda x: x.get('timestamp_ms', 0))

    # Time-based Analysis & Replies
    night_owl_counts = Counter()
    morning_person_counts = Counter()
    replied_to_counts = Counter()

    # Timezone: UTC-8 (PST)
    pst_tz = timezone(timedelta(hours=-8))

    for i, msg in enumerate(all_messages):
        sender = msg.get('sender_name')
        ts_ms = msg.get('timestamp_ms', 0)
        
        # Convert to datetime in PST
        dt = datetime.fromtimestamp(ts_ms / 1000.0, tz=pst_tz)
        hour = dt.hour

        # Night Owl: 12 AM - 5 AM (0 - 5)
        if 0 <= hour < 5:
            night_owl_counts[sender] += 1
        
        # Morning Person: 6 AM - 11 AM (6 - 11)
        elif 6 <= hour < 12:
            morning_person_counts[sender] += 1
        
        # Peak Activity Month
        month_key = dt.strftime('%B %Y')
        month_counts[month_key] += 1
        
        # Most Replied To Heuristic
        if i < len(all_messages) - 1:
            next_msg = all_messages[i+1]
            next_ts = next_msg.get('timestamp_ms', 0)
            next_sender = next_msg.get('sender_name')
            
            if next_sender != sender:
                time_diff = (next_ts - ts_ms) / 1000.0 # seconds
                if time_diff < 300: # 5 minutes
                    replied_to_counts[sender] += 1

    # Average Response Time Analysis
    user_response_times = defaultdict(list)
    for i in range(1, len(all_messages)):
        prev_msg = all_messages[i-1]
        curr_msg = all_messages[i]
        prev_sender = prev_msg.get('sender_name')
        curr_sender = curr_msg.get('sender_name')
        
        if prev_sender != curr_sender:
            time_diff = (curr_msg.get('timestamp_ms', 0) - prev_msg.get('timestamp_ms', 0)) / 1000.0
            # Only count if it's within a reasonable "conversation" window, e.g., 12 hours
            if time_diff < 43200: 
                user_response_times[curr_sender].append(time_diff)

    avg_response_times = []
    for user, times in user_response_times.items():
        if len(times) >= 5:
            avg_response_times.append([user, sum(times) / len(times)])
    
    # Average Message Length
    avg_lengths = []
    for user in user_messages_sent:
        if user_messages_sent[user] > 0:
            avg_lengths.append([user, user_total_chars[user] / user_messages_sent[user]])
    

    # --- Generate Rankings ---

    def get_top_10(counter_obj):
        return counter_obj.most_common(10)

    # 1. Top 10 words per user
    top_words_per_user = {}
    for user, counter in user_word_counts.items():
        top_words_per_user[user] = counter.most_common(10)

    # 2. Top 10 words overall
    top_words_overall = total_word_counts.most_common(10)

    # 3. Top 10 liked messages
    top_liked_messages = sorted(messages_with_likes, key=lambda x: x['likes'], reverse=True)[:10]

    # 6. Top 10 image senders
    top_image_senders = user_image_counts.most_common(10)

    # --- Special Rankings ---
    
    # Night Owls
    top_night_owls = get_top_10(night_owl_counts)
    
    # Morning Person
    top_morning_persons = get_top_10(morning_person_counts)
    
    # Most Replied To
    top_replied_to = get_top_10(replied_to_counts)
    
    # Most Reels Sent
    top_reels_sent = get_top_10(user_reels_sent)
    
    # Longest Word Sent (Top 10 users by longest word length)
    longest_word_list = [{"user": u, "word": d["word"], "length": d["length"]} for u, d in user_longest_words.items()]
    top_longest_words = sorted(longest_word_list, key=lambda x: x['length'], reverse=True)[:10]
    
    # Most Aura (Likes Received / Messages Sent)
    # Higher Ratio = Better Aura (More likes per message)
    aura_scores = []
    for user in user_messages_sent:
        sent = user_messages_sent[user]
        received = user_likes_received[user]
        
        # Filter out low activity users
        if sent < 10:
            continue
            
        # Formula: Received / Sent
        ratio = received / sent
             
        aura_scores.append([user, ratio])
        
    # Sort by ratio DESCENDING (Highest Ratio First)
    top_aura = sorted(aura_scores, key=lambda x: x[1], reverse=True)[:10]
    
    # Most Messages Liked (Likes Given)
    top_messages_liked = get_top_10(user_likes_given)

    # Construct final stats object (Original)
    # Shuffle and pick a diverse set for the quiz
    import random
    random.shuffle(quiz_pool)
    # Try to get unique messages from different senders
    seen_senders = set()
    final_quiz = []
    for msg in quiz_pool:
        if msg['sender'] not in seen_senders or len(seen_senders) > 10:
            final_quiz.append(msg)
            seen_senders.add(msg['sender'])
        if len(final_quiz) >= 15:
            break

    stats = {
        "top_words_per_user": top_words_per_user,
        "top_words_overall": top_words_overall,
        "top_liked_messages": top_liked_messages,
        "top_image_senders": top_image_senders,
        "quiz_pool": final_quiz,
        # Expose detected media basenames and available share methods so the UI can show buttons
        "media_files": sorted(list(collected_media_basenames)),
        "share_story_options": ['html', 'pptx', 'pdf']
    }

    # Construct rankings object (New)
    rankings = {
        "night_owls": top_night_owls,
        "morning_person": top_morning_persons,
        "most_replied_to": top_replied_to,
        "most_reels_sent": top_reels_sent,
        "longest_word_sent": top_longest_words,
        "most_aura": top_aura,
        "most_messages_liked": top_messages_liked,
        "most_profanity": get_top_10(user_profanity_counts),
        "avg_response_times": sorted(avg_response_times, key=lambda x: x[1])[:10], # Fastest first
        "slowest_responders": sorted(avg_response_times, key=lambda x: x[1], reverse=True)[:10],
        "avg_message_lengths": sorted(avg_lengths, key=lambda x: x[1], reverse=True)[:10],
        "peak_month": month_counts.most_common(1)[0] if month_counts else ("None", 0),
    }

    # --- Generate Network Data ---
    participants = sorted(list(user_messages_sent.keys()))
    nodes = [{"id": p} for p in participants]
    links = []
    
    # Calculate combined likes and ratios for all pairs
    pair_data = []
    for i in range(len(participants)):
        for j in range(i + 1, len(participants)):
            p1 = participants[i]
            p2 = participants[j]
            combined_likes = likes_matrix[p1][p2] + likes_matrix[p2][p1]
            combined_likes_given = user_likes_given[p1] + user_likes_given[p2]
            
            if combined_likes > 0 and combined_likes_given > 0:
                ratio = combined_likes / combined_likes_given
                pair_data.append({
                    "source": p1,
                    "target": p2,
                    "value": combined_likes,
                    "ratio": ratio
                })
    
    if pair_data:
        # Sort ratios to find quantile thresholds
        all_ratios = sorted([p['ratio'] for p in pair_data])
        num_links = len(all_ratios)
        
        # Thresholds for 3 even levels (tertiles)
        low_thresh = all_ratios[num_links // 3]
        med_thresh = all_ratios[(2 * num_links) // 3]
        
        for p in pair_data:
            if p['ratio'] <= low_thresh:
                p['strength'] = "Low"
            elif p['ratio'] <= med_thresh:
                p['strength'] = "Medium"
            else:
                p['strength'] = "High"
            links.append(p)

    network_data = {
        "nodes": nodes,
        "links": links
    }

    # --- Chart Data ---
    chart_data = {
        "messages_sent": [{"name": u, "value": v} for u, v in user_messages_sent.most_common(15)],
        "likes_given": [{"name": u, "value": v} for u, v in user_likes_given.most_common(15)],
        "likes_received": [{"name": u, "value": v} for u, v in user_likes_received.most_common(15)],
        "reels_sent": [{"name": u, "value": v} for u, v in user_reels_sent.most_common(15)],
    }

    # --- Share Story Generators ---
    # Helper: find image files in a media root by basename list
    def _find_files_by_basename(basenames, media_root):
        matches = []
        if not media_root or not basenames:
            return matches
        for root, dirs, files in os.walk(media_root):
            for b in basenames:
                if b in files:
                    matches.append(os.path.join(root, b))
        # Keep original order of basenames where possible
        ordered = []
        for b in basenames:
            for p in matches:
                if os.path.basename(p) == b:
                    ordered.append(p)
                    break
        return ordered or matches

    def create_html_deck(image_paths, out_html):
        try:
            import base64
            parts = []
            parts.append('<!doctype html><html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width, initial-scale=1"></head><body style="margin:0;background:#000;">')
            for p in image_paths:
                try:
                    with open(p, 'rb') as f:
                        data = f.read()
                    b64 = base64.b64encode(data).decode('ascii')
                    mime = 'image/png'
                    if re.search(r'\.jpe?g$', p, flags=re.I):
                        mime = 'image/jpeg'
                    parts.append(f'<div style="page-break-after:always;width:100%;height:100vh;display:flex;align-items:center;justify-content:center"><img src="data:{mime};base64,{b64}" style="max-width:100%;max-height:100%;"/></div>')
                except Exception:
                    continue
            parts.append('</body></html>')
            with open(out_html, 'w', encoding='utf-8') as f:
                f.write(''.join(parts))
            return out_html
        except Exception:
            return None

    def create_pptx_from_images(image_paths, out_pptx):
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except Exception:
            return None
        try:
            prs = Presentation()
            # Use blank layout for slides
            blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
            for img in image_paths:
                try:
                    slide = prs.slides.add_slide(blank)
                    left = top = 0
                    slide.shapes.add_picture(img, left, top, width=prs.slide_width)
                except Exception:
                    continue
            prs.save(out_pptx)
            return out_pptx
        except Exception:
            return None

    def create_pdf_from_images(image_paths, out_pdf):
        # Prefer img2pdf for exact image-to-page conversion
        try:
            import img2pdf
            with open(out_pdf, 'wb') as f:
                f.write(img2pdf.convert(image_paths))
            return out_pdf
        except Exception:
            try:
                from PIL import Image
                pil_imgs = []
                for p in image_paths:
                    try:
                        im = Image.open(p)
                        if im.mode == 'RGBA':
                            im = im.convert('RGB')
                        pil_imgs.append(im)
                    except Exception:
                        continue
                if not pil_imgs:
                    return None
                pil_imgs[0].save(out_pdf, save_all=True, append_images=pil_imgs[1:])
                return out_pdf
            except Exception:
                return None

    def generate_share_story(media_root, basenames, out_dir, methods=('html', 'pptx', 'pdf')):
        """
        media_root: path to extracted media files (local dir)
        basenames: list of image basenames to include (from 'media_files')
        out_dir: directory to write outputs
        methods: tuple of methods to attempt
        Returns dict of generated artifact paths.
        """
        os.makedirs(out_dir, exist_ok=True)
        image_paths = _find_files_by_basename(basenames, media_root)
        results = {}
        if not image_paths:
            return results
        if 'html' in methods:
            out_html = os.path.join(out_dir, 'story_deck.html')
            r = create_html_deck(image_paths, out_html)
            if r:
                results['html'] = r
        if 'pptx' in methods:
            out_pptx = os.path.join(out_dir, 'story_deck.pptx')
            r = create_pptx_from_images(image_paths, out_pptx)
            if r:
                results['pptx'] = r
        if 'pdf' in methods:
            out_pdf = os.path.join(out_dir, 'story_deck.pdf')
            r = create_pdf_from_images(image_paths, out_pdf)
            if r:
                results['pdf'] = r
        return results

    # Available share story options for the front-end to offer
    share_story_options = ['html', 'pptx', 'pdf']

    return {
        "stats": stats,
        "rankings": rankings,
        "network_data": network_data,
        "chart_data": chart_data,
        "media_files": sorted(list(collected_media_basenames)),
        "share_story_options": share_story_options,
        # helper function available for local/testing usage; not serializable over JSON
        "_generate_share_story": generate_share_story
    }

if __name__ == "__main__":
    # Simulated file data for testing
    simulated_files_data = None
    
    # For local testing, you can manually set the file data here
    # Example: simulated_files_data = [{"content": "{\"messages\":[{\"sender_name\":\"Alice\",\"content\":\"Hello Bob! 😊\",\"timestamp_ms\":1633072800000}]}}"]
    
    if simulated_files_data:
        results = process_uploaded_data(simulated_files_data)
        # For local execution, you might still want to save the results
        with open('stats.json', 'w', encoding='utf-8') as f:
            json.dump(results['stats'], f, indent=2, ensure_ascii=False)
        print("Stats saved to stats.json")
        with open('rankings.json', 'w', encoding='utf-8') as f:
            json.dump(results['rankings'], f, indent=2, ensure_ascii=False)
        print("Rankings saved to rankings.json")
        with open('network_data.json', 'w', encoding='utf-8') as f:
            json.dump(results['network_data'], f, indent=2, ensure_ascii=False)
        print("Network data saved to network_data.json")
        # Also write share options and media file basenames separately for UI testing
        try:
            with open('share_options.json', 'w', encoding='utf-8') as f:
                json.dump({"share_story_options": results.get('share_story_options', results.get('stats', {}).get('share_story_options', []))}, f, indent=2, ensure_ascii=False)
            print('Share options saved to share_options.json')
        except Exception as e:
            print(f'Error writing share_options.json: {e}')
        
        try:
            with open('media_files.json', 'w', encoding='utf-8') as f:
                json.dump({"media_files": results.get('media_files', results.get('stats', {}).get('media_files', []))}, f, indent=2, ensure_ascii=False)
            print('Media files saved to media_files.json')
        except Exception as e:
            print(f'Error writing media_files.json: {e}')
